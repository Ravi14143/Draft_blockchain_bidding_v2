from flask import Blueprint, jsonify, request, session, current_app, send_from_directory
from werkzeug.utils import secure_filename
from datetime import datetime
from functools import wraps
import os, json 

# Import models
from src.models.user import User, db, RFQ, Bid, Project, Milestone, ClarificationThread, ClarificationMessage, RFQFile
# 👆 added RFQFile import

# Import blockchain service
from src.blockchain.contract_service import (
    create_rfq_onchain, submit_bid_onchain, str_keccak, to_unix_seconds
)

# Import evaluation services (keeping typo as per file name)
from src.services.evalution import evaluate_phase1, evaluate_phase2


user_bp = Blueprint('user', __name__, url_prefix='/api')

# ---------------------------
# Utility decorators
# ---------------------------
def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return jsonify({'error': 'Authentication required'}), 401
        return f(*args, **kwargs)
    return decorated_function

def role_required(role):
    def decorator(f):
        @wraps(f)
        def decorated_function(*args, **kwargs):
            if 'user_id' not in session:
                return jsonify({'error': 'Authentication required'}), 401
            user = User.query.get(session['user_id'])
            if not user or user.role != role:
                return jsonify({'error': 'Insufficient permissions'}), 403
            return f(*args, **kwargs)
        return decorated_function
    return decorator

# ---------------------------
# Authentication routes
# ---------------------------

@user_bp.route('/register', methods=['POST'])
def register_user():
    data = request.json or {}

    username = data.get('username')
    password = data.get('password')
    role = data.get('role')

    if not username or not password or not role:
        return jsonify({'error': 'Username, password, and role are required'}), 400

    if User.query.filter_by(username=username).first():
        return jsonify({'error': 'Username already exists'}), 400

    # Validate bidder fields
    if role == 'bidder':
        required = ['name', 'email', 'phone', 'company', 'address']
        for field in required:
            if not data.get(field):
                return jsonify({'error': f'{field} is required for bidders'}), 400

    # Create user
    user = User(
        username=username,
        role=role,
        name=data.get('name'),
        email=data.get('email'),
        phone=data.get('phone'),
        company=data.get('company'),
        address=data.get('address'),
        avatar_url=data.get('avatar_url')
    )
    user.set_password(password)

    try:
        db.session.add(user)
        db.session.commit()
        # Auto login
        session['user_id'] = user.id
        return jsonify(user.to_dict()), 201
    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f"Registration failed: {str(e)}"}), 500


@user_bp.route('/login', methods=['POST'])
def login():
    data = request.json
    user = User.query.filter_by(username=data['username']).first()
    if user and user.check_password(data['password']):
        session['user_id'] = user.id
        return jsonify(user.to_dict())
    return jsonify({'error': 'Invalid credentials'}), 401


@user_bp.route('/logout', methods=['POST'])
@login_required
def logout_user():
    session.pop('user_id', None)
    return jsonify({'message': 'Logged out successfully'})


@user_bp.route('/me', methods=['GET'])
@login_required
def get_current_user():
    user = User.query.get(session['user_id'])
    print('hi')
    print(user)
    return jsonify(user.to_dict())

# ---------------------------
# Admin routes
# ---------------------------

@user_bp.route('/users', methods=['GET','post'])
@login_required
@role_required('admin')
def get_all_users():
    users = User.query.all()
    return jsonify([user.to_dict() for user in users])

@user_bp.route('/users/<int:user_id>', methods=['DELETE'])
@login_required
@role_required('admin')
def delete_user(user_id):
    user = User.query.get_or_404(user_id)
    db.session.delete(user)
    db.session.commit()
    return '', 204



# ---------------------------
# Get all RFQs
# ---------------------------
@user_bp.route('/rfqs', methods=['GET'])
@login_required
def get_rfqs():
    try:
        rfqs = RFQ.query.all()
        return jsonify([rfq.to_dict(include_files=True) for rfq in rfqs])
    except Exception as e:
        current_app.logger.error(f"Error in get_rfqs: {e}")
        return jsonify({"error": str(e)}), 500


# ---------------------------
# Helper: Save Uploaded File
# ---------------------------
def save_file(file, rfq_id):
    upload_dir = os.path.join(
        current_app.config.get("UPLOAD_FOLDER", "uploads"),
        "rfqs",
        str(rfq_id)
    )
    os.makedirs(upload_dir, exist_ok=True)

    filename = secure_filename(file.filename)
    filepath = os.path.join(upload_dir, filename)
    file.save(filepath)

    rfq_file = RFQFile(rfq_id=rfq_id, filename=filename, filepath=filepath)
    db.session.add(rfq_file)
    return rfq_file


# ---------------------------
# Create RFQ
# ---------------------------
@user_bp.route('/rfqs', methods=['POST'])
@role_required('owner')
def create_rfq():
    try:
        # -------- Parse request --------
        if request.content_type and request.content_type.startswith('multipart/form-data'):
            metadata_str = request.form.get('metadata', '{}')
            data = json.loads(metadata_str)
            files = request.files.getlist('files')
        else:
            data = request.json or {}
            files = []

        # -------- Validate date fields --------
        date_fields = ["deadline", "publish_date", "clarification_deadline", "start_date", "end_date"]
        for field in date_fields:
            if data.get(field):
                try:
                    datetime.strptime(data[field], "%Y-%m-%d")
                except ValueError:
                    raise ValueError(f"Invalid {field} format, expected YYYY-MM-DD")

        deadline_str = data.get("deadline", "")

        # -------- Hash scope --------
        scope_text = data.get('scope', '')
        meta_hash = str_keccak(scope_text)

        # -------- Budget --------
        def parse_int(val):
            try:
                return int(val) if val else None
            except ValueError:
                return None

        budget_min = parse_int(data.get("budget_min"))
        budget_max = parse_int(data.get("budget_max"))
        budget = budget_max or budget_min or 0

        # -------- Normalize evaluation weights --------
        weights_str = data.get("evaluation_weights", "")
        if weights_str and "=" in weights_str:
            # Convert "Price=40,Experience=30" → {"price":40,...}
            weights_dict = {}
            for part in weights_str.split(","):
                if "=" in part:
                    k, v = part.split("=")
                    try:
                        weights_dict[k.strip().lower()] = float(v.strip())
                    except ValueError:
                        continue
            weights_str = json.dumps(weights_dict)

        # -------- On-chain RFQ creation --------
        onchain = create_rfq_onchain(
            data.get('title', ''),
            meta_hash,
            deadline_str,  # YYYY-MM-DD
            data.get('category', ''),
            budget,
            data.get('location', '')
        )

        rfq_id_chain = onchain.get("rfqId")
        tx_hash = onchain.get("txHash")

        if not tx_hash:
            raise ValueError("On-chain RFQ creation failed: no transaction hash.")
        if not rfq_id_chain:
            current_app.logger.warning(
                "On-chain RFQ created but rfqId missing (event parsing failed)."
            )

        # -------- Save RFQ in DB --------
        rfq = RFQ(
            owner_id=session.get('user_id'),
            title=data.get('title', ''),
            scope=scope_text,
            deadline=deadline_str,
            evaluation_criteria=data.get('evaluation_criteria', ''),
            category=data.get('category'),
            budget_min=budget_min,
            budget_max=budget_max,
            publish_date=data.get('publish_date'),
            clarification_deadline=data.get('clarification_deadline'),
            start_date=data.get('start_date'),
            end_date=data.get('end_date'),
            eligibility_requirements=data.get('eligibility_requirements'),
            evaluation_weights=weights_str,
            onchain_id=rfq_id_chain,
            tx_hash=tx_hash,
            status="open"
        )

        db.session.add(rfq)
        db.session.commit()

        # -------- Save uploaded files --------
        for f in files:
            save_file(f, rfq.id)
        db.session.commit()

        return jsonify(rfq.to_dict(include_files=True)), 201

    except Exception as e:
        db.session.rollback()
        current_app.logger.error(f"RFQ creation failed: {e}")
        return jsonify({'error': f"RFQ creation failed: {str(e)}"}), 500


# ---------------------------
# Serve RFQ File
# ---------------------------
from flask import send_file

@user_bp.route('/rfqs/<int:rfq_id>/files/<path:filename>', methods=['GET'])
@login_required
def serve_rfq_file(rfq_id, filename):
    rfq = RFQ.query.get_or_404(rfq_id)

    upload_dir = os.path.join(
        current_app.config.get("UPLOAD_FOLDER", "uploads"),
        "rfqs",
        str(rfq_id)
    )
    file_path = os.path.join(upload_dir, filename)

    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404

    # ✅ This sets correct MIME type and forces download
    return send_file(file_path, as_attachment=True)



# ---------------------------
# Get Single RFQ
# ---------------------------
@user_bp.route('/rfqs/<int:rfq_id>', methods=['GET'])
@login_required
def get_rfq(rfq_id):
    rfq = RFQ.query.get_or_404(rfq_id)
    return jsonify(rfq.to_dict(include_files=True))


# ---------------------------
# Bid routes
# ---------------------------
@user_bp.route('/rfqs/<int:rfq_id>/bids', methods=['GET','post'])
@login_required
def get_rfq_bids(rfq_id):
    user = User.query.get(session['user_id'])
    rfq = RFQ.query.get_or_404(rfq_id)
    if user.role == 'owner' and rfq.owner_id == user.id:
        bids = Bid.query.filter_by(rfq_id=rfq_id).order_by(Bid.created_at.desc()).all()
    elif user.role == 'bidder':
        bids = Bid.query.filter_by(rfq_id=rfq_id, bidder_id=user.id).order_by(Bid.created_at.desc()).all()
    else:
        return jsonify({'error': 'Insufficient permissions'}), 403
    return jsonify([bid.to_dict() for bid in bids])


@user_bp.route('/bids', methods=['POST'])
@role_required('bidder')
def create_bid():
    try:
        if request.content_type.startswith('multipart/form-data'):
            data_str = request.form.get('data', '{}')
            data = json.loads(data_str)
            files = request.files.getlist('files')
        else:
            return jsonify({"error": "Bid must include a file upload"}), 400

        if not files:
            return jsonify({"error": "At least one file (PDF/PPT) is required"}), 400

        # --- Save bid in DB ---
        bid = Bid(
            rfq_id=data['rfq_id'],
            bidder_id=session['user_id'],
            price=data['price'],
            timeline=data.get('timeline'),
            qualifications=data.get('qualifications', ''),
            status="submitted",
            phase1_status="pending",
            phase2_status="pending"
        )
        db.session.add(bid)
        db.session.flush()  # get bid.id

        upload_dir = os.path.join(
            current_app.config.get("UPLOAD_FOLDER", "uploads"),
            "bids",
            str(bid.id)
        )
        os.makedirs(upload_dir, exist_ok=True)

        for f in files:
            filename = secure_filename(f.filename)
            filepath = os.path.join(upload_dir, filename)
            f.save(filepath)

            bid_file = BidFile(bid_id=bid.id, filename=filename, filepath=filepath)
            db.session.add(bid_file)

        db.session.commit()

        # --- Use uploaded file(s) for evaluation ---
        from PyPDF2 import PdfReader
        import pptx

        text_content = ""
        for bf in bid.files:
            if bf.filename.lower().endswith(".pdf"):
                reader = PdfReader(bf.filepath)
                text_content += "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
            elif bf.filename.lower().endswith(".ppt") or bf.filename.lower().endswith(".pptx"):
                prs = pptx.Presentation(bf.filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"

        bid.qualifications = text_content[:5000]  # trim if needed
        db.session.commit()

        # Run evaluations on extracted text
        p1 = evaluate_phase1(bid.qualifications, bid.rfq_id)
        bid.phase1_status = p1.get("status", "pending")
        bid.phase1_report = {
            "reasons": p1.get("reasons", []),
            "missing": p1.get("missing", []),
            "red_flags": p1.get("red_flags", [])
        }
        bid.red_flags = p1.get("red_flags", []) or []

        if bid.phase1_status == "reject":
            bid.status = "rejected"
        elif bid.phase1_status == "clarify":
            bid.status = "needs_clarification"
        else:
            bid.status = "submitted"

        db.session.commit()

        # Phase 2
        if bid.phase1_status == "pass":
            p2 = evaluate_phase2(
                bid_id=bid.id,
                rfq_id=bid.rfq_id,
                qualifications_text=bid.qualifications,
                price=bid.price,
                timeline=bid.timeline
            )
            bid.phase2_status = p2.get("status", "pending")
            bid.phase2_score = p2.get("score")
            bid.phase2_breakdown = p2.get("breakdown")
            bid.red_flags = list(set((bid.red_flags or []) + p2.get("red_flags", [])))

        db.session.commit()

        return jsonify({"bid": bid.to_dict(include_files=True)}), 201

    except Exception as e:
        db.session.rollback()
        return jsonify({'error': f"Bid submission failed: {str(e)}"}), 500


@user_bp.route('/my-bids', methods=['GET','post'])
@role_required('bidder')
def get_my_bids():
    bids = Bid.query.filter_by(bidder_id=session['user_id']).order_by(Bid.created_at.desc()).all()
    return jsonify([bid.to_dict() for bid in bids])

# ---- Owner actions on a bid: reject / clarify ----
@user_bp.route('/bids/<int:bid_id>/reject', methods=['POST'])
@role_required('owner')
def reject_bid(bid_id):
    bid = Bid.query.get_or_404(bid_id)
    rfq = RFQ.query.get_or_404(bid.rfq_id)
    if rfq.owner_id != session['user_id']:
        return jsonify({'error': 'Not authorized'}), 403
    bid.status = "rejected"
    bid.phase1_status = "reject"
    db.session.commit()
    return jsonify({"message": "Bid rejected", "bid": bid.to_dict()})

@user_bp.route('/bids/<int:bid_id>/clarify', methods=['POST'])
@role_required('owner')
def request_clarification(bid_id):
    """
    Owner requests clarification on a bid (Phase 1 -> clarify).
    Body: {"question": "..."}
    """
    bid = Bid.query.get_or_404(bid_id)
    rfq = RFQ.query.get_or_404(bid.rfq_id)
    if rfq.owner_id != session['user_id']:
        return jsonify({'error': 'Not authorized'}), 403

    data = request.json or {}
    question = (data.get("question") or "").strip()
    if not question:
        return jsonify({"error": "Question is required"}), 400

    # open thread
    thread = ClarificationThread(bid_id=bid.id, owner_id=rfq.owner_id, status="open")
    db.session.add(thread)
    db.session.flush()  # get thread.id

    msg = ClarificationMessage(
        thread_id=thread.id, sender_id=rfq.owner_id, role="owner", message=question
    )
    bid.status = "needs_clarification"
    bid.phase1_status = "clarify"
    db.session.add(msg)
    db.session.commit()

    return jsonify({"thread": thread.to_dict(include_messages=True), "bid": bid.to_dict()}), 201

@user_bp.route('/bids/<int:bid_id>/clarifications', methods=['GET','post'])
@login_required
def get_bid_clarifications(bid_id):
    bid = Bid.query.get_or_404(bid_id)
    rfq = RFQ.query.get_or_404(bid.rfq_id)
    user = User.query.get(session['user_id'])
    if not (user.role == 'owner' and rfq.owner_id == user.id) and not (user.role == 'bidder' and bid.bidder_id == user.id):
        return jsonify({'error': 'Not authorized'}), 403
    threads = ClarificationThread.query.filter_by(bid_id=bid.id).order_by(ClarificationThread.created_at.asc()).all()
    return jsonify([t.to_dict(include_messages=True) for t in threads])

@user_bp.route('/clarifications/<int:thread_id>/reply', methods=['POST'])
@login_required
def reply_clarification(thread_id):
    """
    Bidder replies to clarification. Auto re-evaluate Phase 1 (and Phase 2 if passes).
    Body: {"message": "..."}
    """
    thread = ClarificationThread.query.get_or_404(thread_id)
    bid = Bid.query.get_or_404(thread.bid_id)
    rfq = RFQ.query.get_or_404(bid.rfq_id)
    user = User.query.get(session['user_id'])

    # Only bid's bidder or owner can reply (usually bidder)
    if not (user.id == bid.bidder_id or user.id == rfq.owner_id):
        return jsonify({'error': 'Not authorized'}), 403

    data = request.json or {}
    msg_txt = (data.get("message") or "").strip()
    if not msg_txt:
        return jsonify({"error": "Message is required"}), 400

    msg = ClarificationMessage(
        thread_id=thread.id, sender_id=user.id, role=user.role, message=msg_txt
    )
    db.session.add(msg)
    db.session.flush()

    # If bidder replied, re-run Phase 1 using qualifications + entire thread content
    if user.id == bid.bidder_id:
        # concatenate thread messages to feed as context
        all_msgs = ClarificationMessage.query.filter_by(thread_id=thread.id).order_by(ClarificationMessage.created_at.asc()).all()
        context = "\n\n".join([m.role.upper() + ": " + m.message for m in all_msgs])
        combined_text = f"{bid.qualifications}\n\nCLARIFICATION CONTEXT:\n{context}"

        p1 = evaluate_phase1(combined_text, bid.rfq_id)
        bid.phase1_status = p1.get("status", "pending")
        bid.phase1_report = {
            "reasons": p1.get("reasons", []),
            "missing": p1.get("missing", []),
            "red_flags": p1.get("red_flags", [])
        }
        bid.red_flags = list(set((bid.red_flags or []) + p1.get("red_flags", [])))

        if bid.phase1_status == "pass":
            # Run Phase 2 now
            p2 = evaluate_phase2(bid.qualifications + "\n" + msg_txt, bid.price, bid.timeline, bid.rfq_id)
            bid.phase2_status = p2.get("status", "pending")
            bid.phase2_score = p2.get("score")
            bid.phase2_breakdown = p2.get("breakdown")
            bid.red_flags = list(set((bid.red_flags or []) + p2.get("red_flags", [])))

            if bid.phase2_status == "pass":
                bid.status = "submitted"
                thread.status = "resolved"
                thread.resolved_at = datetime.utcnow()
            elif bid.phase2_status == "reject":
                bid.status = "rejected"
                thread.status = "rejected"
                thread.resolved_at = datetime.utcnow()
            else:
                bid.status = "needs_clarification"
        elif bid.phase1_status == "reject":
            bid.status = "rejected"
            thread.status = "rejected"
            thread.resolved_at = datetime.utcnow()
        else:
            bid.status = "needs_clarification"

    db.session.commit()
    return jsonify({
        "thread": thread.to_dict(include_messages=True),
        "bid": bid.to_dict()
    })

# Winner selection
@user_bp.route('/bids/<int:bid_id>/select', methods=['POST'])
@role_required('owner')
def select_winner(bid_id):
    bid = Bid.query.get_or_404(bid_id)
    rfq = RFQ.query.get_or_404(bid.rfq_id)
    if rfq.owner_id != session['user_id']:
        return jsonify({'error': 'Not authorized'}), 403
    project = Project(rfq_id=rfq.id, winner_bid_id=bid.id)
    rfq.status = 'closed'
    bid.status = 'selected'
    db.session.add(project)
    db.session.commit()
    return jsonify(project.to_dict()), 201

# ---- Candidates view for owner ----
@user_bp.route('/rfqs/<int:rfq_id>/candidates', methods=['GET','post'])
@role_required('owner')
def get_candidates(rfq_id):
    rfq = RFQ.query.get_or_404(rfq_id)
    if rfq.owner_id != session['user_id']:
        return jsonify({'error': 'Not authorized'}), 403

    bids = Bid.query.filter_by(rfq_id=rfq.id)\
        .filter(Bid.phase1_status == "pass")\
        .filter(Bid.phase2_status == "pass")\
        .order_by(Bid.phase2_score.desc())\
        .all()

    return jsonify([{
        **b.to_dict(),
        "candidate": True
    } for b in bids])

# ---------------------------
# Project & Milestone routes
# ---------------------------
@user_bp.route('/projects', methods=['GET','POST'])
@login_required
def get_projects():
    user = User.query.get(session['user_id'])
    if user.role == 'owner':
        projects = Project.query.join(RFQ).filter(RFQ.owner_id == user.id).all()
    elif user.role == 'bidder':
        projects = Project.query.join(Bid).filter(Bid.bidder_id == user.id).all()
    else:
        projects = Project.query.all()
    return jsonify([project.to_dict() for project in projects])

@user_bp.route('/projects/<int:project_id>', methods=['GET','post'])
@login_required
def get_project(project_id):
    project = Project.query.get_or_404(project_id)
    return jsonify(project.to_dict())

@user_bp.route('/projects/<int:project_id>/milestones', methods=['GET','post'])
@login_required
def get_project_milestones(project_id):
    milestones = Milestone.query.filter_by(project_id=project_id).all()
    return jsonify([milestone.to_dict() for milestone in milestones])

@user_bp.route('/milestones', methods=['POST'])
@role_required('bidder')
def create_milestone():
    data = request.json
    milestone = Milestone(
        project_id=data['project_id'],
        description=data['description'],
        due_date=data['due_date'],
        document_hash=data.get('document_hash', 'hash_placeholder')
    )
    db.session.add(milestone)
    db.session.commit()
    return jsonify(milestone.to_dict()), 201

@user_bp.route('/milestones/<int:milestone_id>/approve', methods=['POST'])
@role_required('owner')
def approve_milestone(milestone_id):
    milestone = Milestone.query.get_or_404(milestone_id)
    milestone.status = 'approved'
    db.session.commit()
    return jsonify(milestone.to_dict())


# ---------------------------
# Dashboard
# ---------------------------
@user_bp.route('/dashboard', methods=['GET','post'])
@login_required
def dashboard():
    user = User.query.get(session['user_id'])
    if not user:
        return jsonify({"error": "User not found"}), 404

    # ---------------- Bidder Dashboard ----------------
    if user.role == "bidder":
        available_rfqs = RFQ.query.filter_by(status="open").order_by(RFQ.created_at.desc()).all()
        bids = Bid.query.filter_by(bidder_id=user.id).order_by(Bid.created_at.desc()).all()
        projects = Project.query.join(Bid).filter(Bid.bidder_id == user.id, Bid.status == "selected").all()

        return jsonify({
            "role": "bidder",
            "available_rfqs": len(available_rfqs),
            "bid_count": len(bids),
            "project_count": len(projects),
            "recent_rfqs": [r.to_dict() for r in available_rfqs[:5]],
            "bids": [b.to_dict() for b in bids[:5]]
        })

    # ---------------- Owner Dashboard ----------------
    if user.role == "owner":
        rfqs = RFQ.query.filter_by(owner_id=user.id).order_by(RFQ.created_at.desc()).all()
        projects = Project.query.join(RFQ).filter(RFQ.owner_id == user.id).all()
        return jsonify({
            "role": "owner",
            "rfq_count": len(rfqs),
            "project_count": len(projects),
            "rfqs": [r.to_dict() for r in rfqs],
            "projects": [p.to_dict() for p in projects]
        })

    # ---------------- Admin Dashboard ----------------
    if user.role == "admin":
        users = User.query.all()
        return jsonify({
            "role": "admin",
            "user_count": len(users),
            "users": [u.to_dict() for u in users]
        })

    return jsonify({"error": "Unknown role"}), 400



# -----------------------------
# GET bidder profile
# -----------------------------
@user_bp.route('/bidder/profile', methods=['GET'])
@role_required('bidder')
def get_bidder_profile():
    user = User.query.get(session['user_id'])
    if not user:
        return jsonify({"error": "User not found"}), 404
    return jsonify(user.to_dict()), 200

# -----------------------------
# PUT bidder profile
# -----------------------------
@user_bp.route('/bidder/profile', methods=['PUT'])
@role_required('bidder')
def update_bidder_profile():
    user = User.query.get(session['user_id'])
    if not user:
        return jsonify({"error": "User not found"}), 404

    data = request.json or {}

    # Update only allowed fields
    user.name = data.get("name", user.name)
    user.email = data.get("email", user.email)
    user.phone = data.get("phone", user.phone)
    user.company = data.get("company", user.company)
    user.address = data.get("address", user.address)
    user.avatar_url = data.get("avatar_url", user.avatar_url)

    try:
        db.session.commit()
        return jsonify({"message": "Profile updated successfully"}), 200
    except Exception as e:
        db.session.rollback()
        return jsonify({"error": f"Failed to update profile: {str(e)}"}), 500